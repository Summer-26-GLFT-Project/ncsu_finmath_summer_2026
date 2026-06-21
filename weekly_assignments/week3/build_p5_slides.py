"""Build a two-slide, one-minute PS3 Problem 5 presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


TEMPLATE = Path(r"C:\Users\Ruleigh\Downloads\ncsu_template.pptx")
HERE = Path(__file__).resolve().parent
OUTPUT = HERE / "ps3_p5_price_impact_slides.pptx"
CHART = HERE / "figures" / "ps3_p5_mean_pnl_degradation.png"

RED = RGBColor(204, 0, 0)
DARK_RED = RGBColor(153, 0, 0)
BLACK = RGBColor(25, 25, 25)
DARK_GRAY = RGBColor(70, 70, 70)
MID_GRAY = RGBColor(220, 220, 220)
LIGHT_GRAY = RGBColor(245, 245, 245)
WHITE = RGBColor(255, 255, 255)


def delete_all_slides(prs: Presentation) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def add_text(slide, text, x, y, w, h, *, size=28, color=BLACK,
             bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.08)
    frame.margin_top = frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rounded_box(slide, x, y, w, h, *, fill=LIGHT_GRAY, line=MID_GRAY,
                radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(
        radius_shape, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    return shape


def arrow(slide, x, y, w=0.65, h=0.42):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.color.rgb = RED
    return shape


def set_title(slide, text):
    title = slide.shapes.title
    title.text = text
    for p in title.text_frame.paragraphs:
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = BLACK


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


prs = Presentation(TEMPLATE)
delete_all_slides(prs)

# ---------------------------------------------------------------------------
# Slide 1: mechanism
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
set_title(slide, "What happens after a fill?")

add_text(slide, "BID FILL", 1.0, 1.55, 2.2, 0.55, size=24, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
tag = rounded_box(slide, 1.0, 1.48, 2.2, 0.70, fill=RED, line=RED)
slide.shapes._spTree.remove(tag._element)
slide.shapes._spTree.insert(2, tag._element)

steps_y = 2.65
for x, heading, detail in [
    (0.75, "We buy 1 unit", "inventory q increases"),
    (4.10, "Reservation price falls", "r = S − qγσ²τ"),
    (7.45, "Both quotes move down", "buy less · sell sooner"),
]:
    rounded_box(slide, x, steps_y, 2.75, 1.35)
    add_text(slide, heading, x + 0.15, steps_y + 0.17, 2.45, 0.43,
             size=20, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, detail, x + 0.15, steps_y + 0.68, 2.45, 0.34,
             size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)
arrow(slide, 3.55, steps_y + 0.47)
arrow(slide, 6.90, steps_y + 0.47)

add_text(slide, "ASK FILL", 16.95, 1.55, 2.2, 0.55, size=24, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
tag2 = rounded_box(slide, 16.95, 1.48, 2.2, 0.70, fill=DARK_RED, line=DARK_RED)
slide.shapes._spTree.remove(tag2._element)
slide.shapes._spTree.insert(2, tag2._element)

for x, heading, detail in [
    (10.75, "We sell 1 unit", "inventory q decreases"),
    (14.10, "Reservation price rises", "move back toward q = 0"),
    (17.15, "Both quotes move up", "buy sooner · sell less"),
]:
    rounded_box(slide, x, steps_y, 2.45, 1.35)
    add_text(slide, heading, x + 0.12, steps_y + 0.17, 2.21, 0.43,
             size=19, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, detail, x + 0.12, steps_y + 0.68, 2.21, 0.34,
             size=15, color=DARK_GRAY, align=PP_ALIGN.CENTER)
arrow(slide, 13.52, steps_y + 0.47, w=0.48)
arrow(slide, 16.67, steps_y + 0.47, w=0.40)

rounded_box(slide, 1.0, 5.10, 18.0, 2.15, fill=WHITE, line=RED)
add_text(slide, "Two different effects", 1.35, 5.35, 4.2, 0.55,
         size=24, color=RED, bold=True)
add_text(slide, "A–S quote response", 5.45, 5.30, 3.65, 0.48,
         size=20, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "$0.40 per fill", 5.45, 5.88, 3.65, 0.62,
         size=30, color=RED, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "≠", 9.40, 5.62, 1.0, 0.65, size=34, color=DARK_GRAY,
         bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Adverse market impact tested", 10.65, 5.30, 4.5, 0.48,
         size=20, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "$0.05 / $0.10 per fill", 10.65, 5.88, 4.5, 0.62,
         size=30, color=RED, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Goal: test how much the independence assumption overstates P&L",
         1.0, 8.75, 18.0, 0.75, size=26, color=BLACK, bold=True,
         align=PP_ALIGN.CENTER)
add_notes(slide, "A fill means one of our posted orders trades. A bid fill makes us long, so the A–S model moves both quotes down to reduce inventory; an ask fill does the reverse. Do not confuse the forty-cent quote response with the five- or ten-cent adverse market impact tested in Problem 5.")

# ---------------------------------------------------------------------------
# Slide 2: result
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_title(slide, "Price impact erodes expected P&L")
picture = slide.shapes.add_picture(str(CHART), Inches(0.65), Inches(1.72),
                                   width=Inches(12.2), height=Inches(7.30))
picture.crop_top = 0.07

add_text(slide, "1,000 simulated days", 13.35, 1.55, 5.4, 0.5,
         size=19, color=DARK_GRAY, align=PP_ALIGN.CENTER)

for y, eta, pnl, loss in [
    (2.30, "$0.00 / fill", "$2,573", "baseline"),
    (4.20, "$0.05 / fill", "$2,476", "−3.8%"),
    (6.10, "$0.10 / fill", "$2,379", "−7.5%"),
]:
    rounded_box(slide, 13.35, y, 5.4, 1.48, fill=LIGHT_GRAY,
                line=RED if loss != "baseline" else MID_GRAY)
    add_text(slide, eta, 13.65, y + 0.17, 2.35, 0.38,
             size=18, color=DARK_GRAY, bold=True)
    add_text(slide, pnl, 13.65, y + 0.62, 2.35, 0.57,
             size=30, color=BLACK, bold=True)
    add_text(slide, loss, 16.25, y + 0.48, 2.05, 0.55,
             size=27, color=RED if loss != "baseline" else DARK_GRAY,
             bold=True, align=PP_ALIGN.RIGHT)

rounded_box(slide, 13.35, 8.15, 5.4, 1.27, fill=RED, line=RED)
add_text(slide, "Ignoring own-price impact\noverstates profitability",
         13.60, 8.30, 4.90, 0.82, size=22, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_notes(slide, "Across one thousand simulated days, mean terminal P&L falls from about 2,573 dollars with no impact to 2,476 dollars at five cents per fill and 2,379 dollars at ten cents per fill. That is a 3.8 and 7.5 percent reduction. The nearly straight line follows from our linear impact assumption: some apparent spread profit disappears once fills predict adverse price moves.")

prs.save(OUTPUT)
print(OUTPUT)
